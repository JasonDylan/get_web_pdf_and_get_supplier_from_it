# This code takes all pdfs in documents, scans them, then generates paragraph embeddings on a sliding scale of 200 tokens (roughly 150 words)
# Run this before you run EmbedDocuments.py or app.py
# You need an OpenAI key saved in APIkey.txt
# Note that if your PDFs are not searchable, this won't work - use a third party tool to convert them to txt or doc first.  You
#   can look at the "-originaltext.csv" file created here and scan real quick to see if the text looks corrupted for any of your docs


import os
import time
import chardet
from PyPDF2 import PdfReader 
import nltk
import pandas as pd
import numpy as np
import json
import io
import re
from pptx import Presentation
# you need to pip install python-docx, not docx
import docx
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')

# Set the desired chunk size and overlap size
# chunk_size is how many tokens we will take in each block of text
# overlap_size is how much overlap. So 200, 100 gives you chunks of between the 1st and 200th word, the 100th and 300th, the 200 and 400th...
# I have in no way optimized these
chunk_size = 2000
overlap_size = 100

# load user settings and api key
def read_settings(file_name):
    settings = {}
    with open(file_name, "r") as f:
        for line in f:
            key, value = line.strip().split("=")
            settings[key] = value
    return settings
settings = read_settings("settings.txt")
filedirectory = "data/report_pdf"
# Check if the subfolder exists, if not, create it
output_folder = "Textchunks"
if not os.path.exists(output_folder):
    os.makedirs(output_folder)

# 遍历文件夹及其子文件夹中的所有文件
def list_files_recursive(directory):
    file_list = []
    for root, _, files in os.walk(directory):
        for file_path in files:
            file_path = os.path.join(root, file_path)
            file_list.append(file_path)
    return file_list

# 从指定文件夹及其子文件夹中获取所有文件列表
file_list = list_files_recursive(filedirectory)

# 打印所有文件的路径
for file_path in file_list:
    print(file_path)
# Loop through all pdf, txt, tex in the "documents" folder
for file_path in file_list:
    # Create an empty DataFrame to store the text and title of each document
    df = pd.DataFrame(columns=["Title", "Text"])
    print("Loading " + file_path)
    if file_path.endswith(".pdf"):
        # Open the PDF file in read-binary mode
        reader = PdfReader(file_path)

        # Extract the text from each page of the PDF
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"

        # Add the text and title to the DataFrame
        title = os.path.splitext(os.path.basename(file_path))[0]  # Remove the file extension from the file_path
        new_row = pd.DataFrame({"Title": [title], "Text": [text]})
        df = pd.concat([df, new_row], ignore_index=True)

    elif file_path.endswith(".ppt") or file_path.endswith(".pptx"):
        ppt = Presentation(file_path)
        text = ''
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
        title = os.path.splitext(file_path)[0]  # Remove the file extension from the file_path
        new_row = pd.DataFrame({"Title": [title], "Text": [text]})
        df = pd.concat([df, new_row], ignore_index=True)

    elif file_path.endswith(".doc") or file_path.endswith(".docx"):
        # Open the DOC/DOCX file in binary mode and read the raw data
        doc = docx.Document(file_path)

        # Convert the file to UTF-8 and extract the text
        text = ''
        for paragraph in doc.paragraphs:
            text += paragraph.text

        # Add the text and title to the DataFrame
        title = os.path.splitext(file_path)[0]  # Remove the file extension from the file_path
        new_row = pd.DataFrame({"Title": [title], "Text": [text]})
        df = pd.concat([df, new_row], ignore_index=True)

    elif file_path.endswith(".txt"):
        # Open the text file and read its contents
        with open(file_path, "r", encoding="utf-8") as file:
            text = file.read()

        # Add the text and title to the DataFrame
        title = os.path.splitext(file_path)[0]  # Remove the file extension from the file_path
        new_row = pd.DataFrame({"Title": [title], "Text": [text]})
        df = pd.concat([df, new_row], ignore_index=True)
        
    elif file_path.endswith(".tex"):
        # Use regular expressions to extract regular text from the LaTeX file
        with open(file_path, "r", encoding="utf-8") as file:
            text = file.read()

        # Replace special characters
        text = text.replace('\\$', '$')
        text = text.replace('\\\\', '\n')  # Replace \\ with newline for paragraph breaks
        # Remove comments
        text = re.sub(r'%.*?\n', '', text)

        def replace_math_expression(match):
            # Remove $ or $$ signs but keep the expression
            return match.group(1)

        # Modified regular expression to match both $...$ and $$...$$
        text= re.sub(r'\${1,2}(.*?)\${1,2}', replace_math_expression, text)

        # Remove \begin{} ... \end{} blocks
        text = re.sub(r'\\begin{.*?}.*?\\end{.*?}', '', text, flags=re.DOTALL)

        # Remove common LaTeX commands
        commands = [
            r'\\textbf{.*?}', r'\\textit{.*?}', r'\\emph{.*?}', r'\\underline{.*?}',  # Formatting
            r'\\cite{.*?}', r'\\ref{.*?}',  # References
            r'\\label{.*?}',  # Labels
            # Add more commands as needed
        ]
        for command in commands:
            text = re.sub(command, '', text)
        
        # Add the text and title to the DataFrame
        title = os.path.splitext(file_path)[0] # Remove the file extension from the file_path
        new_row = pd.DataFrame({"Title": [title], "Text": [text]})
        df = pd.concat([df, new_row], ignore_index=True)

    # Loop through the rows and create overlapping chunks for each text
    chunks = []
    for i, row in df.iterrows():
        # Tokenize the text for the current row
        tokens = nltk.word_tokenize(row['Text'])

        # Loop through the tokens and create overlapping chunks
        for j in range(0, len(tokens), chunk_size - overlap_size):
            # Get the start and end indices of the current chunk
            start = j
            end = j + chunk_size

            # Create the current chunk by joining the tokens within the start and end indices
            chunk = ' '.join(tokens[start:end])

            # Add the article title to the beginning of the chunk
            chunk_with_title = "This text comes from the document " + row['Title'] + ". " + chunk

            # Append the current chunk to the list of chunks, along with the corresponding title
            chunks.append([row['Title'], chunk_with_title])

    # Convert the list of chunks to a dataframe
    df_chunks = pd.DataFrame(chunks, columns=['Title', 'Text'])

    # Truncate the file_path if it's too long, e.g., limit to 250 characters
    max_file_path_length = 250
    if len(file_path) > max_file_path_length:
        file_path = file_path[:max_file_path_length]

    # Remove the file extension from the file_path
    file_path_without_extension = os.path.splitext(os.path.basename(file_path))[0]

    # 构建输出文件夹路径
    company_name = file_path.replace(filedirectory, "").replace(os.path.basename(file_path),"").strip("\\") # 获取相对路径
    output_subfolder = os.path.join(output_folder, company_name)
    os.makedirs(output_subfolder, exist_ok=True)  # 创建子文件夹
    # Save the df_chunks to the output_folder subfolder with the new file name
    output_file = os.path.join(output_subfolder, file_path_without_extension + "-originaltext.csv")
    df_chunks.to_csv(output_file, encoding='utf-8', escapechar='\\', index=False)

    print("Saving " + file_path)
